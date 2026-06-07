from fastapi import FastAPI, Request, Response, HTTPException
from pptx import Presentation
import io
import requests
import os

app = FastAPI()

@app.get("/")
def health():
    return {"status": "ok", "service": "pptx-render", "version": "1.0"}

@app.post("/render")
async def render_pptx(request: Request):
    try:
        data = await request.json()
        template_url = data.get('template_url')
        placeholders = data.get('placeholders', {})
        
        if not template_url:
            raise HTTPException(status_code=400, detail="template_url is required")
        
        # Fetch template
        r = requests.get(template_url, timeout=30)
        if r.status_code != 200:
            raise HTTPException(status_code=400, detail=f"Failed to fetch template: HTTP {r.status_code}")
        
        # Load presentation
        prs = Presentation(io.BytesIO(r.content))
        
        # Replace placeholders in all slides
        total_replacements = 0
        replacement_log = {}
        
        for slide in prs.slides:
            slide_log = replace_in_shapes(slide.shapes, placeholders)
            for key, count in slide_log.items():
                replacement_log[key] = replacement_log.get(key, 0) + count
                total_replacements += count
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Filename
        date_filename = placeholders.get('DATE_FILENAME', 'output')
        filename = f"Daily_Market_Update_{date_filename}.pptx"
        
        return Response(
            content=output.read(),
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            headers={
                'Content-Disposition': f'attachment; filename="{filename}"',
                'X-Total-Replacements': str(total_replacements),
                'X-Filename': filename
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"{type(e).__name__}: {str(e)}")


def replace_in_shapes(shapes, placeholders):
    """Recursively replace placeholders in shapes (text frames, tables, groups)."""
    log = {}
    for shape in shapes:
        # Text frames
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                _replace_in_paragraph(para, placeholders, log)
        
        # Tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        _replace_in_paragraph(para, placeholders, log)
        
        # Grouped shapes
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            sub_log = replace_in_shapes(shape.shapes, placeholders)
            for k, v in sub_log.items():
                log[k] = log.get(k, 0) + v
    
    return log


def _replace_in_paragraph(para, placeholders, log):
    """Replace placeholders in a paragraph, handling cross-run splits."""
    full_text = para.text
    
    # Check if any placeholder is present
    has_placeholder = any(
        f'{{{{{key}}}}}' in full_text 
        for key in placeholders.keys() 
        if not key.startswith('_')
    )
    
    if not has_placeholder:
        return
    
    # Try per-run replacement first (preserves formatting per run)
    pending_replacements = {}
    for key, value in placeholders.items():
        if key.startswith('_'):
            continue
        ph = f'{{{{{key}}}}}'
        if ph not in full_text:
            continue
        
        # Try to find in single run
        found = False
        for run in para.runs:
            if ph in run.text:
                count = run.text.count(ph)
                run.text = run.text.replace(ph, str(value or ''))
                log[key] = log.get(key, 0) + count
                found = True
                break
        
        if not found:
            pending_replacements[key] = value
    
    # Fallback: consolidate runs if placeholders split across them
    if pending_replacements:
        new_text = para.text
        for key, value in pending_replacements.items():
            ph = f'{{{{{key}}}}}'
            count = new_text.count(ph)
            new_text = new_text.replace(ph, str(value or ''))
            log[key] = log.get(key, 0) + count
        
        runs = list(para.runs)
        if runs:
            runs[0].text = new_text
            for r in runs[1:]:
                r.text = ''


if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8080))
    uvicorn.run(app, host="0.0.0.0", port=port)
